#!/usr/bin/env python3
"""
Knowledge Base Ingestion Script — Python 3.9 compatible
Converts Documentacion/ into .discovery/knowledge/
Tags: AS-IS (legacy/source technology) | TO-BE (target technology)
"""

from __future__ import annotations

import argparse
import hashlib
import json
import os
import re
from datetime import datetime
from pathlib import Path
from typing import Optional

# Defaults: relative to current working directory
DEFAULT_INPUT_BASE = Path("input")
DEFAULT_KB_BASE = Path(".discovery/knowledge")
INPUT_BASE = DEFAULT_INPUT_BASE
KB_BASE = DEFAULT_KB_BASE

# Map top-level folder names to classification.
# Add project-specific entries via --folder-tag 'Name=as-is:Tech Label'
FOLDER_TAGS: dict[str, dict[str, str]] = {}

LEGACY_HINTS = ("legacy", "old", "source", "origin", "v3", "v4")
TARGET_HINTS = ("target", "destino", "new", "v5", "v6")

SKIP_EXTENSIONS = {".eot", ".ttf", ".woff", ".woff2", ".ico", ".js", ".css"}
SKIP_SVG_PATTERNS = ["roboto-v", "ionicons"]
SKIP_SVG_NAMES = {"coverage-badge-documentation.svg", "compodoc-vectorise.svg", "compodoc-vectorise-inverted.svg"}
SKIP_PNG_NAMES = {"compodoc-vectorise.png", "compodoc-vectorise-inverted.png"}

INVENTORY = []
ERRORS = []
NOW = datetime.now().isoformat()


# ─── Helpers ───────────────────────────────────────────────────────────────────


def slugify(name: str) -> str:
    name = os.path.splitext(name)[0]
    name = re.sub(r"[^\w\- ()]", "", name, flags=re.UNICODE)
    return name.strip()


def tag_for_path(rel_path: str) -> dict:
    parts = Path(rel_path).parts
    top = parts[0] if parts else "."
    if len(parts) == 1:
        return {
            "classification": "to-be",
            "technology": "Shared documentation",
            "source_root": ".",
        }
    if top in FOLDER_TAGS:
        return FOLDER_TAGS[top] | {"source_root": top}

    lowered = top.lower()
    if any(hint in lowered for hint in LEGACY_HINTS):
        return {
            "classification": "as-is",
            "technology": top,
            "source_root": top,
        }
    if any(hint in lowered for hint in TARGET_HINTS):
        return {
            "classification": "to-be",
            "technology": top,
            "source_root": top,
        }

    return {
        "classification": "to-be",
        "technology": top,
        "source_root": top,
    }


def output_dir_for(rel_path: str) -> Path:
    parts = Path(rel_path).parts
    tag_info = tag_for_path(rel_path)
    d = KB_BASE / "ingested" / tag_info["classification"]
    if len(parts) > 1:
        d = d / parts[0]
        for p in parts[1:-1]:
            d = d / p
    d.mkdir(parents=True, exist_ok=True)
    return d


def ingested_dir() -> Path:
    d = KB_BASE / "ingested"
    d.mkdir(parents=True, exist_ok=True)
    return d


def artifact_stem(rel_path: str) -> str:
    stem = str(Path(rel_path).with_suffix(""))
    stem = stem.replace(os.sep, "--")
    stem = re.sub(r"[^\w\s-]", "", stem, flags=re.UNICODE)
    stem = re.sub(r"[\s_]+", "-", stem)
    stem = re.sub(r"-{2,}", "-", stem)
    return stem.strip("-").lower() or "document"


def file_sha256(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as handle:
        for chunk in iter(lambda: handle.read(65536), b""):
            digest.update(chunk)
    return digest.hexdigest()


def build_frontmatter(source: str, tag_info: dict, extra: Optional[dict] = None) -> str:
    lines = [
        "---",
        f"source: {source}",
        f"tag: {tag_info['classification']}",
        f"classification: {tag_info['classification']}",
        f"category: {tag_info['classification']}",
        f"technology: {tag_info['technology']}",
        f"source_root: {tag_info['source_root']}",
        f"converted_at: {NOW}",
        "converter: knowledge-analysis",
    ]
    if extra:
        for k, v in extra.items():
            lines.append(f"{k}: {v}")
    lines.append("---\n")
    return "\n".join(lines)


def build_registry(input_base: Path, kb_base: Path) -> dict:
    documents = {}
    by_classification = {}

    for entry in INVENTORY:
        rel_path = entry["relative"]
        outputs = []
        for output in entry.get("outputs", []):
            try:
                outputs.append(str(Path(output).resolve().relative_to(kb_base.resolve())))
            except ValueError:
                outputs.append(output)

        doc_entry = dict(entry)
        doc_entry["tag"] = entry["classification"]
        doc_entry["outputs"] = outputs
        documents[rel_path] = doc_entry
        by_classification.setdefault(entry["classification"], []).append(rel_path)

    return {
        "generated_at": NOW,
        "input_base": str(input_base),
        "output_base": str(kb_base),
        "documents": documents,
        "by_classification": {
            classification: sorted(paths) for classification, paths in sorted(by_classification.items())
        },
    }


def build_state(input_base: Path, kb_base: Path, stats: dict) -> dict:
    documents = {}
    file_hashes = {}

    for entry in INVENTORY:
        rel_path = entry["relative"]
        outputs = []
        for output in entry.get("outputs", []):
            try:
                outputs.append(str(Path(output).resolve().relative_to(kb_base.resolve())))
            except ValueError:
                outputs.append(output)

        documents[rel_path] = {
            "classification": entry["classification"],
            "technology": entry["technology"],
            "format": entry["format"],
            "outputs": outputs,
            "converted_at": entry["converted_at"],
            "sha256": entry.get("sha256"),
            "size_bytes": entry.get("size_bytes"),
        }
        if entry.get("sha256"):
            file_hashes[rel_path] = entry["sha256"]

    return {
        "status": "completed_with_errors" if ERRORS else "completed",
        "updated_at": NOW,
        "input_base": str(input_base),
        "output_base": str(kb_base),
        "ingested_dir": str(kb_base / "ingested"),
        "stats": stats,
        "file_hashes": file_hashes,
        "documents": documents,
        "errors": ERRORS,
    }


def load_existing_state(kb_base: Path) -> dict:
    state_path = kb_base / "state.json"
    if not state_path.exists():
        return {}
    try:
        return json.loads(state_path.read_text(encoding="utf-8"))
    except Exception:
        return {}


def outputs_exist(kb_base: Path, outputs: list[str]) -> bool:
    if not outputs:
        return False
    for output in outputs:
        output_path = Path(output)
        if not output_path.is_absolute():
            output_path = kb_base / output_path
        if not output_path.exists():
            return False
    return True


# ─── PDF ───────────────────────────────────────────────────────────────────────


def convert_pdf(input_path: Path, rel_path: str) -> list:
    import pdfplumber

    tag_info = tag_for_path(rel_path)
    out_dir = output_dir_for(rel_path)
    stem = slugify(input_path.name)
    md_lines = []
    all_tables = []

    with pdfplumber.open(str(input_path)) as pdf:
        n_pages = len(pdf.pages)
        md_lines.append(build_frontmatter(str(input_path), tag_info, {"pages": n_pages, "format": "PDF"}))
        md_lines.append(f"# {stem}\n")

        for i, page in enumerate(pdf.pages):
            page_text = page.extract_text() or ""
            page_tables = page.extract_tables() or []

            if page_text.strip():
                md_lines.append(f"<!-- Page {i + 1} -->")
                md_lines.append(page_text.strip())
                md_lines.append("")

            for j, table in enumerate(page_tables):
                if not table:
                    continue
                headers = [str(c or "") for c in table[0]]
                md_lines.append(f"\n**Table (page {i + 1}, #{j + 1})**\n")
                md_lines.append("| " + " | ".join(headers) + " |")
                md_lines.append("| " + " | ".join(["---"] * len(headers)) + " |")
                for row in table[1:]:
                    cells = [str(c or "") for c in row]
                    md_lines.append("| " + " | ".join(cells) + " |")
                md_lines.append("")
                all_tables.append(
                    {
                        "page": i + 1,
                        "table": j + 1,
                        "headers": headers,
                        "rows": [[str(c or "") for c in r] for r in table[1:]],
                    }
                )

    md_path = out_dir / f"{stem}.md"
    md_path.write_text("\n".join(md_lines), encoding="utf-8")
    outputs = [str(md_path)]

    if all_tables:
        json_path = out_dir / f"{stem}.tables.json"
        json_path.write_text(
            json.dumps({"source": str(input_path), "tables": all_tables}, indent=2, ensure_ascii=False),
            encoding="utf-8",
        )
        outputs.append(str(json_path))

    INVENTORY.append(
        {
            "source": str(input_path),
            "relative": rel_path,
            "classification": tag_info["classification"],
            "technology": tag_info["technology"],
            "source_root": tag_info["source_root"],
            "format": "PDF",
            "pages": n_pages,
            "outputs": outputs,
            "converted_at": NOW,
        }
    )
    return outputs


# ─── DOCX ──────────────────────────────────────────────────────────────────────


def convert_docx(input_path: Path, rel_path: str) -> list:
    import docx as docx_lib

    tag_info = tag_for_path(rel_path)
    out_dir = output_dir_for(rel_path)
    stem = slugify(input_path.name)

    doc = docx_lib.Document(str(input_path))
    props = doc.core_properties

    extra = {"format": "DOCX", "paragraphs": len(doc.paragraphs), "tables": len(doc.tables)}
    if props.title:
        extra["title"] = props.title
    if props.author:
        extra["author"] = props.author

    heading_map = {
        "Heading 1": "#",
        "Heading 2": "##",
        "Heading 3": "###",
        "Heading 4": "####",
        "Heading 5": "#####",
        "Title": "#",
    }

    md_lines = [build_frontmatter(str(input_path), tag_info, extra)]
    md_lines.append(f"# {props.title or stem}\n")

    for para in doc.paragraphs:
        style = para.style.name if para.style else ""
        text = para.text.strip()
        if not text:
            md_lines.append("")
            continue
        if style in heading_map:
            md_lines.append(f"{heading_map[style]} {text}\n")
        elif "List Bullet" in style:
            md_lines.append(f"- {text}")
        elif "List Number" in style:
            md_lines.append(f"1. {text}")
        else:
            runs_md = ""
            for run in para.runs:
                t = run.text
                if run.bold and run.italic:
                    runs_md += f"***{t}***"
                elif run.bold:
                    runs_md += f"**{t}**"
                elif run.italic:
                    runs_md += f"*{t}*"
                else:
                    runs_md += t
            md_lines.append(runs_md if runs_md.strip() else text)

    all_tables = []
    for i, table in enumerate(doc.tables):
        md_lines.append(f"\n### Table {i + 1}\n")
        rows = []
        for j, row in enumerate(table.rows):
            cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
            rows.append(cells)
            if j == 0:
                md_lines.append("| " + " | ".join(cells) + " |")
                md_lines.append("| " + " | ".join(["---"] * len(cells)) + " |")
            else:
                md_lines.append("| " + " | ".join(cells) + " |")
        all_tables.append(
            {
                "table_index": i,
                "headers": rows[0] if rows else [],
                "rows": rows[1:] if rows else [],
            }
        )

    md_path = out_dir / f"{stem}.md"
    md_path.write_text("\n".join(md_lines), encoding="utf-8")
    outputs = [str(md_path)]

    if all_tables:
        json_path = out_dir / f"{stem}.tables.json"
        json_path.write_text(
            json.dumps({"source": str(input_path), "tables": all_tables}, indent=2, ensure_ascii=False),
            encoding="utf-8",
        )
        outputs.append(str(json_path))

    INVENTORY.append(
        {
            "source": str(input_path),
            "relative": rel_path,
            "classification": tag_info["classification"],
            "technology": tag_info["technology"],
            "source_root": tag_info["source_root"],
            "format": "DOCX",
            "outputs": outputs,
            "converted_at": NOW,
        }
    )
    return outputs


# ─── PPTX ──────────────────────────────────────────────────────────────────────


def convert_pptx(input_path: Path, rel_path: str) -> list:
    import io
    import tempfile
    import zipfile

    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    tag_info = tag_for_path(rel_path)
    out_dir = output_dir_for(rel_path)
    stem = slugify(input_path.name)

    actual_path = str(input_path)
    tmp_path = None

    if input_path.suffix.lower() == ".potx":
        tmp_file = tempfile.NamedTemporaryFile(suffix=".pptx", delete=False)
        tmp_path = tmp_file.name
        tmp_file.close()

        buffer = io.BytesIO()
        with zipfile.ZipFile(str(input_path), "r") as src_zip:
            with zipfile.ZipFile(buffer, "w", zipfile.ZIP_DEFLATED) as dst_zip:
                for item in src_zip.infolist():
                    data = src_zip.read(item.filename)
                    if item.filename == "[Content_Types].xml":
                        data = data.replace(
                            b"presentationml.template.main+xml",
                            b"presentationml.presentation.main+xml",
                        )
                    dst_zip.writestr(item, data)

        with open(tmp_path, "wb") as patched_file:
            patched_file.write(buffer.getvalue())
        actual_path = tmp_path

    try:
        prs = Presentation(actual_path)
    finally:
        if tmp_path and os.path.exists(tmp_path):
            os.unlink(tmp_path)

    n_slides = len(prs.slides)

    md_lines = [build_frontmatter(str(input_path), tag_info, {"format": "PPTX", "slides": n_slides})]
    md_lines.append(f"# {stem}\n")

    slides_data = []
    image_count = 0

    for i, slide in enumerate(prs.slides):
        layout_name = slide.slide_layout.name if slide.slide_layout else "unknown"
        slide_info = {"slide_number": i + 1, "layout": layout_name, "shapes": [], "notes": ""}

        title_text = ""
        if slide.shapes.title and slide.shapes.title.text:
            title_text = slide.shapes.title.text.strip()

        header = f"## Slide {i + 1}"
        if title_text:
            header += f": {title_text}"
        md_lines.append(header + "\n")

        for shape in slide.shapes:
            if shape.has_text_frame:
                texts = []
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if not text:
                        continue
                    level = para.level or 0
                    indent = "  " * level
                    line = f"{indent}- {text}" if level > 0 else text
                    texts.append(line)
                if texts:
                    md_lines.extend(texts)
                    md_lines.append("")
                    slide_info["shapes"].append({"type": "text", "content": "\n".join(texts)})

            elif shape.has_table:
                table = shape.table
                all_rows = list(table.rows)  # _RowCollection doesn't support slicing
                headers = [c.text.strip() for c in all_rows[0].cells] if all_rows else []
                if headers:
                    md_lines.append("| " + " | ".join(headers) + " |")
                    md_lines.append("| " + " | ".join(["---"] * len(headers)) + " |")
                rows_data = []
                for row in all_rows[1:]:
                    cells = [c.text.strip() for c in row.cells]
                    md_lines.append("| " + " | ".join(cells) + " |")
                    rows_data.append(cells)
                md_lines.append("")
                slide_info["shapes"].append({"type": "table", "headers": headers, "rows": rows_data})

            elif hasattr(shape, "shape_type") and shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image_count += 1
                md_lines.append(f"*[Image {image_count}]*\n")
                slide_info["shapes"].append({"type": "image", "index": image_count})

        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                md_lines.append(f"> **Notes:** {notes}\n")
                slide_info["notes"] = notes

        md_lines.append("---\n")
        slides_data.append(slide_info)

    md_path = out_dir / f"{stem}.md"
    md_path.write_text("\n".join(md_lines), encoding="utf-8")

    json_path = out_dir / f"{stem}.json"
    json_path.write_text(
        json.dumps({"source": str(input_path), "slides": slides_data}, indent=2, ensure_ascii=False),
        encoding="utf-8",
    )

    outputs = [str(md_path), str(json_path)]
    INVENTORY.append(
        {
            "source": str(input_path),
            "relative": rel_path,
            "classification": tag_info["classification"],
            "technology": tag_info["technology"],
            "source_root": tag_info["source_root"],
            "format": "PPTX",
            "slides": n_slides,
            "images": image_count,
            "outputs": outputs,
            "converted_at": NOW,
        }
    )
    return outputs


# ─── XLSX ──────────────────────────────────────────────────────────────────────


def convert_xlsx(input_path: Path, rel_path: str) -> list:
    import openpyxl

    tag_info = tag_for_path(rel_path)
    out_dir = output_dir_for(rel_path)
    stem = slugify(input_path.name)

    wb = openpyxl.load_workbook(str(input_path), data_only=False)
    md_lines = [build_frontmatter(str(input_path), tag_info, {"format": "XLSX", "sheets": len(wb.sheetnames)})]
    md_lines.append(f"# {stem}\n")
    sheets_data = []

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows = []
        headers = []
        formulas = []

        for i, row in enumerate(ws.iter_rows(values_only=False)):
            cell_values = []
            for cell in row:
                val = cell.value
                if isinstance(val, str) and val.startswith("="):
                    formulas.append({"cell": cell.coordinate, "formula": val})
                cell_values.append(str(val) if val is not None else "")
            if i == 0:
                headers = cell_values
            else:
                if any(v.strip() for v in cell_values):
                    rows.append(cell_values)

        sheets_data.append(
            {
                "name": sheet_name,
                "headers": headers,
                "row_count": len(rows),
                "rows": rows,
                "formulas": formulas or None,
            }
        )

        md_lines.append(f"## {sheet_name}\n")
        extras = f"Rows: {len(rows)} | Columns: {len(headers)}"
        if formulas:
            extras += f" | Formulas: {len(formulas)}"
        md_lines.append(extras + "\n")

        if headers:
            cols = min(len(headers), 10)
            md_lines.append("| " + " | ".join(str(h) for h in headers[:cols]) + " |")
            md_lines.append("| " + " | ".join(["---"] * cols) + " |")
            for row in rows[:20]:
                md_lines.append("| " + " | ".join(str(v) for v in row[:cols]) + " |")
            if len(rows) > 20:
                md_lines.append(f"\n*... {len(rows) - 20} more rows*\n")
        md_lines.append("")

    json_path = out_dir / f"{stem}.json"
    json_path.write_text(
        json.dumps({"source": str(input_path), "sheets": sheets_data}, indent=2, ensure_ascii=False, default=str),
        encoding="utf-8",
    )

    md_path = out_dir / f"{stem}.md"
    md_path.write_text("\n".join(md_lines), encoding="utf-8")

    outputs = [str(json_path), str(md_path)]
    INVENTORY.append(
        {
            "source": str(input_path),
            "relative": rel_path,
            "classification": tag_info["classification"],
            "technology": tag_info["technology"],
            "source_root": tag_info["source_root"],
            "format": "XLSX",
            "sheets": list(wb.sheetnames),
            "outputs": outputs,
            "converted_at": NOW,
        }
    )
    return outputs


# ─── HTML ──────────────────────────────────────────────────────────────────────


def convert_html(input_path: Path, rel_path: str) -> list:
    from bs4 import BeautifulSoup

    tag_info = tag_for_path(rel_path)
    out_dir = output_dir_for(rel_path)
    stem = slugify(input_path.name)

    html = input_path.read_text(encoding="utf-8", errors="replace")
    soup = BeautifulSoup(html, "html.parser")

    for tag in soup(["script", "style", "nav", "footer", "noscript", "header"]):
        tag.decompose()

    page_title = soup.title.get_text(strip=True) if soup.title else stem

    main = soup.find("div", class_="content") or soup.find("main") or soup.body or soup

    md_lines = [
        build_frontmatter(str(input_path), tag_info, {"format": "HTML", "title": page_title}),
        f"# {page_title}\n",
    ]

    def process_element(el):
        if el.name in ("h1", "h2", "h3", "h4", "h5", "h6"):
            level = int(el.name[1])
            text = el.get_text(" ", strip=True)
            if text:
                md_lines.append(f"\n{'#' * level} {text}\n")
        elif el.name == "p":
            text = el.get_text(" ", strip=True)
            if text:
                md_lines.append(text + "\n")
        elif el.name in ("ul", "ol"):
            for li in el.find_all("li", recursive=False):
                text = li.get_text(" ", strip=True)
                if text:
                    md_lines.append(f"- {text}")
            md_lines.append("")
        elif el.name == "table":
            rows_data = el.find_all("tr")
            if not rows_data:
                return
            header_cells = rows_data[0].find_all(["th", "td"])
            headers = [c.get_text(" ", strip=True) for c in header_cells]
            if headers:
                md_lines.append("| " + " | ".join(headers) + " |")
                md_lines.append("| " + " | ".join(["---"] * len(headers)) + " |")
                for row in rows_data[1:]:
                    cells = [c.get_text(" ", strip=True) for c in row.find_all(["td", "th"])]
                    while len(cells) < len(headers):
                        cells.append("")
                    md_lines.append("| " + " | ".join(cells[: len(headers)]) + " |")
                md_lines.append("")
        elif el.name in ("pre", "code"):
            text = el.get_text()
            if text.strip():
                md_lines.append(f"\n```\n{text.strip()}\n```\n")
        elif el.name and el.name not in ("script", "style", "nav", "footer"):
            for child in el.children:
                if hasattr(child, "name"):
                    process_element(child)

    if main:
        for child in main.children:
            if hasattr(child, "name"):
                process_element(child)

    # Collapse multiple blank lines
    result_lines = []
    prev_blank = False
    for line in md_lines:
        is_blank = line.strip() == ""
        if is_blank and prev_blank:
            continue
        result_lines.append(line)
        prev_blank = is_blank

    md_path = out_dir / f"{stem}.md"
    md_path.write_text("\n".join(result_lines), encoding="utf-8")

    outputs = [str(md_path)]
    INVENTORY.append(
        {
            "source": str(input_path),
            "relative": rel_path,
            "classification": tag_info["classification"],
            "technology": tag_info["technology"],
            "source_root": tag_info["source_root"],
            "format": "HTML",
            "title": page_title,
            "outputs": outputs,
            "converted_at": NOW,
        }
    )
    return outputs


# ─── SVG ───────────────────────────────────────────────────────────────────────


def convert_svg(input_path: Path, rel_path: str) -> list:
    from bs4 import BeautifulSoup

    tag_info = tag_for_path(rel_path)
    out_dir = output_dir_for(rel_path)
    stem = slugify(input_path.name)

    svg_text = input_path.read_text(encoding="utf-8", errors="replace")
    soup = BeautifulSoup(svg_text, "xml")
    texts = [t.get_text(strip=True) for t in soup.find_all("text") if t.get_text(strip=True)]

    md_lines = [
        build_frontmatter(str(input_path), tag_info, {"format": "SVG"}),
        f"# {stem}\n",
        f"*SVG diagram — {len(texts)} text nodes extracted*\n",
    ]
    if texts:
        md_lines.append("## Text content\n")
        for t in texts:
            md_lines.append(f"- {t}")

    md_path = out_dir / f"{stem}.md"
    md_path.write_text("\n".join(md_lines), encoding="utf-8")

    outputs = [str(md_path)]
    INVENTORY.append(
        {
            "source": str(input_path),
            "relative": rel_path,
            "classification": tag_info["classification"],
            "technology": tag_info["technology"],
            "source_root": tag_info["source_root"],
            "format": "SVG",
            "text_nodes": len(texts),
            "outputs": outputs,
            "converted_at": NOW,
        }
    )
    return outputs


def convert_markdown(input_path: Path, rel_path: str) -> list:
    tag_info = tag_for_path(rel_path)
    out_dir = output_dir_for(rel_path)
    stem = artifact_stem(rel_path)

    content = input_path.read_text(encoding="utf-8", errors="replace").strip()
    lines = content.splitlines()
    heading_pattern = re.compile(r"^(#{1,6})\s+(.*\S)\s*$")
    link_pattern = re.compile(r"\[([^\]]+)\]\(([^)]+)\)")

    headings = []
    links = []
    title = input_path.stem

    for line in lines:
        match = heading_pattern.match(line)
        if match:
            level = len(match.group(1))
            text = match.group(2).strip()
            headings.append({"level": level, "text": text})
            if level == 1 and title == input_path.stem:
                title = text
        for label, target in link_pattern.findall(line):
            links.append({"label": label, "target": target})

    md_path = out_dir / f"{stem}.md"
    md_path.write_text(
        build_frontmatter(
            str(input_path),
            tag_info,
            {
                "format": "MD",
                "relative_path": rel_path,
                "line_count": len(lines),
                "heading_count": len(headings),
                "link_count": len(links),
            },
        )
        + (content + "\n" if content else ""),
        encoding="utf-8",
    )

    metadata = {
        "source": str(input_path),
        "relative": rel_path,
        "classification": tag_info["classification"],
        "technology": tag_info["technology"],
        "source_root": tag_info["source_root"],
        "format": "MD",
        "title": title,
        "line_count": len(lines),
        "heading_count": len(headings),
        "headings": headings,
        "link_count": len(links),
        "links": links,
        "converted_at": NOW,
        "converter": "knowledge-analysis",
        "output_md": str(md_path),
    }

    json_path = out_dir / f"{stem}.json"
    json_path.write_text(
        json.dumps(metadata, indent=2, ensure_ascii=False),
        encoding="utf-8",
    )

    outputs = [str(md_path), str(json_path)]
    INVENTORY.append(
        {
            "source": str(input_path),
            "relative": rel_path,
            "classification": tag_info["classification"],
            "technology": tag_info["technology"],
            "source_root": tag_info["source_root"],
            "format": "MD",
            "title": title,
            "outputs": outputs,
            "converted_at": NOW,
        }
    )
    return outputs


# ─── Dispatcher ────────────────────────────────────────────────────────────────


def should_skip(rel_path: str, ext: str):
    if ext in SKIP_EXTENSIONS:
        return True, f"skipped extension {ext}"
    name = Path(rel_path).name.lower()
    if ext == ".svg":
        if name in SKIP_SVG_NAMES:
            return True, f"asset SVG ({name})"
        for pat in SKIP_SVG_PATTERNS:
            if pat in name:
                return True, f"font SVG ({pat})"
    if ext == ".png" and name in SKIP_PNG_NAMES:
        return True, "compodoc logo asset"
    return False, ""


def dispatch(input_path: Path, rel_path: str):
    ext = input_path.suffix.lower()
    try:
        if ext == ".pdf":
            return "PDF", convert_pdf(input_path, rel_path)
        elif ext == ".docx":
            return "DOCX", convert_docx(input_path, rel_path)
        elif ext in (".pptx", ".potx"):
            return "PPTX", convert_pptx(input_path, rel_path)
        elif ext in (".xlsx", ".xls"):
            return "XLSX", convert_xlsx(input_path, rel_path)
        elif ext in (".html", ".htm"):
            return "HTML", convert_html(input_path, rel_path)
        elif ext in (".md", ".markdown"):
            return "MD", convert_markdown(input_path, rel_path)
        elif ext == ".svg":
            return "SVG", convert_svg(input_path, rel_path)
        else:
            return "SKIP", []
    except Exception as e:
        import traceback

        ERRORS.append({"file": rel_path, "error": str(e), "trace": traceback.format_exc()})
        return "ERROR", []


# ─── Main ──────────────────────────────────────────────────────────────────────


def main():
    global INPUT_BASE
    global KB_BASE
    global NOW

    parser = argparse.ArgumentParser(description="Ingest documentation into .discovery/knowledge")
    parser.add_argument(
        "input_dir",
        nargs="?",
        default=str(DEFAULT_INPUT_BASE),
        help="Documentation directory to ingest (default: ./input)",
    )
    parser.add_argument(
        "--output-dir",
        default=str(DEFAULT_KB_BASE),
        help="Knowledge base output directory (default: .discovery/knowledge)",
    )
    parser.add_argument(
        "--folder-tag",
        action="append",
        default=[],
        metavar="'Name=classification:Tech Label'",
        help="Map a top-level folder name to a classification. E.g. 'AngularJS=as-is:AngularJS legacy'",
    )
    args = parser.parse_args()

    # Parse --folder-tag entries
    for ft in args.folder_tag:
        if "=" in ft:
            name, rest = ft.split("=", 1)
            if ":" in rest:
                cls, tech = rest.split(":", 1)
            else:
                cls, tech = rest, name
            FOLDER_TAGS[name.strip()] = {
                "classification": cls.strip(),
                "technology": tech.strip(),
            }

    INPUT_BASE = Path(args.input_dir).expanduser().resolve()
    KB_BASE = Path(args.output_dir).expanduser().resolve()
    NOW = datetime.now().isoformat()
    INVENTORY.clear()
    ERRORS.clear()

    print("📚 Knowledge Base Ingestion")
    print(f"   Input:  {INPUT_BASE}")
    print(f"   Output: {KB_BASE}\n")

    KB_BASE.mkdir(parents=True, exist_ok=True)
    (KB_BASE / "ingested" / "as-is").mkdir(parents=True, exist_ok=True)
    (KB_BASE / "ingested" / "to-be").mkdir(parents=True, exist_ok=True)
    existing_state = load_existing_state(KB_BASE)
    existing_docs = existing_state.get("documents", {})

    all_files = []
    for root, dirs, files in os.walk(INPUT_BASE):
        dirs[:] = sorted([d for d in dirs if not d.startswith(".")])
        for f in sorted(files):
            if f.startswith(".") or f.startswith("~$"):
                continue
            full = Path(root) / f
            rel = os.path.relpath(str(full), str(INPUT_BASE))
            all_files.append((full, rel))

    stats = {
        "processed": 0,
        "skipped": 0,
        "already_ingested": 0,
        "errors": 0,
        "by_format": {},
        "by_classification": {},
    }

    for full_path, rel_path in all_files:
        ext = full_path.suffix.lower()
        skip, reason = should_skip(rel_path, ext)
        if skip:
            print(f"  ⏭  SKIP  {rel_path}  ({reason})")
            stats["skipped"] += 1
            continue

        sha256 = file_sha256(full_path)
        existing_doc = existing_docs.get(rel_path, {})
        if existing_doc.get("sha256") == sha256 and outputs_exist(KB_BASE, existing_doc.get("outputs", [])):
            tag_info = tag_for_path(rel_path)
            doc_format = existing_doc.get("format", ext.lstrip(".").upper())
            classification = existing_doc.get("classification", tag_info["classification"])
            INVENTORY.append(
                {
                    "source": str(full_path),
                    "relative": rel_path,
                    "classification": classification,
                    "technology": existing_doc.get("technology", tag_info["technology"]),
                    "source_root": tag_info["source_root"],
                    "format": doc_format,
                    "outputs": existing_doc.get("outputs", []),
                    "converted_at": existing_doc.get("converted_at", NOW),
                    "sha256": sha256,
                    "size_bytes": full_path.stat().st_size,
                }
            )
            print(f"  ⏭  SKIP  {rel_path}  (already ingested)")
            stats["skipped"] += 1
            stats["already_ingested"] += 1
            stats["by_format"][doc_format] = stats["by_format"].get(doc_format, 0) + 1
            stats["by_classification"][classification] = stats["by_classification"].get(classification, 0) + 1
            continue

        fmt, outputs = dispatch(full_path, rel_path)
        if fmt == "SKIP":
            print(f"  ⏭  SKIP  {rel_path}")
            stats["skipped"] += 1
        elif fmt == "ERROR":
            err = ERRORS[-1]["error"] if ERRORS else "unknown"
            print(f"  ❌ ERROR {rel_path}: {err}")
            stats["errors"] += 1
        else:
            classification = tag_for_path(rel_path)["classification"]
            if INVENTORY and INVENTORY[-1]["relative"] == rel_path:
                INVENTORY[-1]["sha256"] = sha256
                INVENTORY[-1]["size_bytes"] = full_path.stat().st_size
            print(f"  ✅ {fmt:5} [{classification:6}] {rel_path}")
            for o in outputs:
                print(f"         → {os.path.relpath(o, str(KB_BASE))}")
            stats["processed"] += 1
            stats["by_format"][fmt] = stats["by_format"].get(fmt, 0) + 1
            stats["by_classification"][classification] = stats["by_classification"].get(classification, 0) + 1

    index = {
        "generated_at": NOW,
        "input_base": str(INPUT_BASE),
        "output_base": str(KB_BASE),
        "stats": stats,
        "documents": INVENTORY,
        "errors": ERRORS,
    }
    index_path = KB_BASE / "index.json"
    index_path.write_text(json.dumps(index, indent=2, ensure_ascii=False), encoding="utf-8")

    registry = build_registry(INPUT_BASE, KB_BASE)
    registry_path = KB_BASE / "registry.json"
    registry_path.write_text(
        json.dumps(registry, indent=2, ensure_ascii=False),
        encoding="utf-8",
    )

    state = build_state(INPUT_BASE, KB_BASE, stats)
    state_path = KB_BASE / "state.json"
    state_path.write_text(
        json.dumps(state, indent=2, ensure_ascii=False),
        encoding="utf-8",
    )

    print(f"\n{'=' * 60}")
    print(f"✅ Processed: {stats['processed']} files")
    print(f"⏭  Skipped:   {stats['skipped']} files (fonts, assets, already ingested, etc.)")
    print(f"   ↳ Already ingested: {stats['already_ingested']}")
    print(f"❌ Errors:    {stats['errors']} files")
    print("\nBy format:")
    for fmt, count in sorted(stats["by_format"].items()):
        print(f"  {fmt}: {count}")
    print("\nBy classification:")
    for classification, count in sorted(stats["by_classification"].items()):
        print(f"  {classification}: {count}")
    print(f"\n📄 Inventory: {index_path}")
    print(f"🗂️  Registry:  {registry_path}")
    print(f"🧭 State:     {state_path}")

    if ERRORS:
        print("\n⚠️  Errors:")
        for e in ERRORS:
            print(f"  {e['file']}: {e['error']}")


if __name__ == "__main__":
    main()
